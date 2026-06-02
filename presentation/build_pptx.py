"""
Build the Maison Cafu slide deck (.pptx) — editorial, grid-based, SF Pro Display.

Design system
-------------
- 16:9 canvas, 12-column grid with consistent gutters and margins.
- Modular type scale, generous vertical rhythm, restrained palette.
- Typeface: SF Pro Display (headings) + SF Pro Text (body). Falls back gracefully
  if the viewer doesn't have them installed.
- Warm-charcoal background, terracotta accent, cream text — same tokens as the app.

Run:
    python presentation/build_pptx.py
Output:
    presentation/assets/Maison_Cafu_Presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x0E, 0x0C)   # near-black warm charcoal
PANEL    = RGBColor(0x1A, 0x16, 0x10)   # card fill
PANEL_HI = RGBColor(0x20, 0x1B, 0x14)   # raised card
TERRA    = RGBColor(0xE1, 0x78, 0x2D)   # accent
TERRA_DK = RGBColor(0xB5, 0x5C, 0x1E)
CREAM    = RGBColor(0xF7, 0xF3, 0xED)   # primary text
SOFT     = RGBColor(0xCE, 0xC6, 0xBA)   # body text
MUTED    = RGBColor(0x8C, 0x84, 0x79)   # secondary
FAINT    = RGBColor(0x1E, 0x1A, 0x14)   # watermark numerals
HAIR     = RGBColor(0x33, 0x2E, 0x29)   # hairlines / borders

# ── type ──────────────────────────────────────────────────────────────────────
DISPLAY = "SF Pro Display"
TEXT    = "SF Pro Text"

# ── canvas + grid ─────────────────────────────────────────────────────────────
EMU_W = Emu(12192000)   # 13.333 in
EMU_H = Emu(6858000)    # 7.5 in

MARGIN_X = 838200       # 0.92 in
MARGIN_TOP = 620000
GUTTER = 152400         # 0.167 in
COLS = 12
CONTENT_W = int(EMU_W) - 2 * MARGIN_X
COL_W = (CONTENT_W - (COLS - 1) * GUTTER) / COLS


def colx(i):
    """Left edge (EMU) of grid column i (0-based)."""
    return int(MARGIN_X + i * (COL_W + GUTTER))


def span_w(n):
    """Width (EMU) spanning n columns."""
    return int(n * COL_W + (n - 1) * GUTTER)


prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ── low-level helpers ─────────────────────────────────────────────────────────
def _no_shadow(shape):
    shape.shadow.inherit = False


def _tracking(run, hundredths_pt):
    """Letter-spacing in 1/100 pt (negative tightens)."""
    run.font._rPr.set("spc", str(int(hundredths_pt)))


def fill_bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    _no_shadow(r)
    tree = slide.shapes._spTree
    tree.remove(r._element)
    tree.insert(2, r._element)
    return r


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    s = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    _no_shadow(s)
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def run(p, text, size, color, *, font=TEXT, bold=False, italic=False, tracking=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    if tracking is not None:
        _tracking(r, tracking)
    return r


def para(tf, first=False, space_before=0, space_after=0, line=1.0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    p.alignment = align
    return p


# ── shared chrome ─────────────────────────────────────────────────────────────
FOOTER_Y = 6420000
INDEX_TOTAL = 14


def header(slide, kicker, index):
    # kicker, top-left
    tf = textbox(slide, colx(0), MARGIN_TOP, span_w(9), 300000)
    p = para(tf, first=True)
    run(p, kicker.upper(), 12.5, TERRA, font=DISPLAY, bold=True, tracking=240)
    # index, top-right
    tf = textbox(slide, colx(9), MARGIN_TOP, span_w(3), 300000)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, f"{index:02d}", 12.5, CREAM, font=DISPLAY, bold=True, tracking=120)
    run(p, f"  /  {INDEX_TOTAL:02d}", 12.5, MUTED, font=DISPLAY, tracking=120)


def title_block(slide, title, y=1000000, size=37, rule=True):
    tf = textbox(slide, colx(0), y, span_w(11), 900000)
    p = para(tf, first=True, line=1.02)
    run(p, title, size, CREAM, font=DISPLAY, bold=True, tracking=-30)
    if rule:
        rect(slide, colx(0), y + 770000, 660000, Pt(4), fill=TERRA)


def footer(slide):
    rect(slide, colx(0), FOOTER_Y, CONTENT_W, Pt(0.75), fill=HAIR)
    tf = textbox(slide, colx(0), FOOTER_Y + 70000, span_w(8), 260000)
    p = para(tf, first=True)
    run(p, "Maison Cafu — SRWA", 9.5, MUTED, font=TEXT, tracking=60)
    tf = textbox(slide, colx(7), FOOTER_Y + 70000, span_w(5), 260000)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, "CEC418 · Software Construction", 9.5, MUTED, font=TEXT, tracking=60)


def watermark(slide, text):
    tf = textbox(slide, colx(7), 3650000, span_w(5), 2600000, anchor=MSO_ANCHOR.BOTTOM)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT, line=0.9)
    run(p, text, 300, FAINT, font=DISPLAY, bold=True, tracking=-200)


# ── content renderers ─────────────────────────────────────────────────────────
CONTENT_TOP = 2050000


def lead(slide, text, y=CONTENT_TOP, cols=10):
    tf = textbox(slide, colx(0), y, span_w(cols), 700000)
    p = para(tf, first=True, line=1.18)
    run(p, text, 18.5, SOFT, font=TEXT)
    return tf


def card(slide, x, y, w, h, label, body):
    rect(slide, x, y, w, h, fill=PANEL, line=HAIR, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
    pad = 200000
    tf = textbox(slide, x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    p = para(tf, first=True, space_after=6, line=1.05)
    run(p, label, 16.5, TERRA, font=DISPLAY, bold=True, tracking=10)
    p = para(tf, line=1.16)
    run(p, body, 14.5, SOFT, font=TEXT)


def card_grid(slide, items, top, *, height=1520000, vgap=200000):
    """2-column grid of labelled cards (max 4 looks best)."""
    hgap = GUTTER * 2
    cw = (CONTENT_W - hgap) / 2
    xs = [colx(0), int(colx(0) + cw + hgap)]
    for idx, (label, body) in enumerate(items):
        col = idx % 2
        row = idx // 2
        y = top + row * (height + vgap)
        card(slide, xs[col], y, cw, height, label, body)


def numbered_rows(slide, items, top, *, row_h=560000):
    """Full-width rows: big terracotta numeral + label + body."""
    for i, (label, body) in enumerate(items):
        y = top + i * row_h
        # numeral
        tf = textbox(slide, colx(0), y, span_w(1), row_h, anchor=MSO_ANCHOR.TOP)
        p = para(tf, first=True, line=1.0)
        run(p, f"{i + 1:02d}", 22, TERRA, font=DISPLAY, bold=True, tracking=-20)
        # label + body
        tf = textbox(slide, colx(1), y + 20000, span_w(11), row_h)
        p = para(tf, first=True, line=1.1)
        run(p, label + "   ", 17, CREAM, font=DISPLAY, bold=True)
        run(p, body, 15.5, SOFT, font=TEXT)
        if i < len(items) - 1:
            rect(slide, colx(1), y + row_h - 70000, span_w(11), Pt(0.75), fill=HAIR)


# ── slide templates ───────────────────────────────────────────────────────────
def s_content(kicker, title, lead_text, cards, index, *, wm=None):
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide)
    header(slide, kicker, index)
    title_block(slide, title)
    y = CONTENT_TOP
    if lead_text:
        lead(slide, lead_text, y=y)
        y += 680000
    card_grid(slide, cards, y)
    footer(slide)
    return slide


def s_rows(kicker, title, lead_text, rows, index):
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide)
    header(slide, kicker, index)
    title_block(slide, title)
    y = CONTENT_TOP
    if lead_text:
        lead(slide, lead_text, y=y)
        y += 720000
    numbered_rows(slide, rows, y)
    footer(slide)
    return slide


# ── SLIDE 1 — cover ───────────────────────────────────────────────────────────
def cover():
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide)
    # faint oversize wordmark watermark
    tf = textbox(slide, colx(6), 3200000, span_w(7), 3200000, anchor=MSO_ANCHOR.BOTTOM)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT, line=0.86)
    run(p, "SRWA", 230, FAINT, font=DISPLAY, bold=True, tracking=-200)

    rect(slide, colx(0), 2520000, 760000, Pt(5), fill=TERRA)

    tf = textbox(slide, colx(0), 2680000, span_w(11), 1500000)
    p = para(tf, first=True, line=0.98)
    run(p, "Maison Cafu", 66, CREAM, font=DISPLAY, bold=True, tracking=-120)
    p = para(tf, space_before=8, line=1.1)
    run(p, "A smart restaurant web app — and the craft of building it well.",
        20, MUTED, font=TEXT, tracking=0)

    # metadata row, bottom, grid-aligned
    meta = [
        ("PROJECT", "Smart Restaurant Web App"),
        ("COURSE", "CEC418 — Software Construction"),
        ("FOCUS", "5 fundamentals · Git · Docker · Jenkins"),
    ]
    for i, (k, v) in enumerate(meta):
        x = colx(i * 4)
        rect(slide, x, 5180000, Pt(2.5), 360000, fill=TERRA)
        tf = textbox(slide, x + 110000, 5180000, span_w(4) - 110000, 380000)
        p = para(tf, first=True, space_after=3)
        run(p, k, 10.5, TERRA, font=DISPLAY, bold=True, tracking=200)
        p = para(tf, line=1.05)
        run(p, v, 14, CREAM, font=TEXT)
    footer(slide)


# ── SLIDE 14 — closing ────────────────────────────────────────────────────────
def closing():
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide)
    header(slide, "Finale · the 60-second story", 14)
    rect(slide, colx(0), 1700000, 760000, Pt(5), fill=TERRA)
    tf = textbox(slide, colx(0), 2050000, span_w(11), 3400000, anchor=MSO_ANCHOR.TOP)
    p = para(tf, first=True, line=1.22)
    run(p, "“Maison Cafu is a React ", 25, CREAM, font=DISPLAY, italic=True)
    run(p, "frontend", 25, TERRA, font=DISPLAY, bold=True, italic=True)
    run(p, ", a PHP ", 25, CREAM, font=DISPLAY, italic=True)
    run(p, "backend", 25, TERRA, font=DISPLAY, bold=True, italic=True)
    run(p, ", and a MySQL ", 25, CREAM, font=DISPLAY, italic=True)
    run(p, "database", 25, TERRA, font=DISPLAY, bold=True, italic=True)
    run(p, ". It minimizes complexity, anticipates change, is built for verification, "
           "reuses design tokens and a shared API, and follows standards like PSR-12. "
           "Git tracks every change, GitHub is the shared home, Docker makes it run "
           "anywhere, and Jenkins tests and builds it on every change.”",
        25, CREAM, font=DISPLAY, italic=True)
    tf = textbox(slide, colx(0), 5450000, span_w(11), 360000)
    p = para(tf, first=True)
    run(p, "Five fundamentals", 14, TERRA, font=DISPLAY, bold=True, tracking=20)
    run(p, "   ·   four tools   ·   one clear story.", 14, MUTED, font=TEXT)
    footer(slide)


# ──────────────────────────────────────────────────────────────────────────────
#  Build the deck
# ──────────────────────────────────────────────────────────────────────────────
cover()

s_content(
    "Chapter 1 · the product", "What is Maison Cafu?",
    "A website for a Cameroonian restaurant. No app to download, no account — a visitor "
    "can do four simple things.",
    [
        ("Browse", "The menu — real dishes (Ndolé, Poulet DG, Mbongo Tchobi) with photos, prices and tags."),
        ("Recommend", "Add a dish and the site suggests what pairs well — like a waiter's advice."),
        ("Order", "Pick dishes, set a quantity, send the order through."),
        ("Reserve", "Book a table for a chosen date and time."),
    ],
    index=2,
)

s_content(
    "Chapter 2 · the architecture", "The three moving parts",
    "Every website has three pieces. Picture a real restaurant.",
    [
        ("Dining room → Frontend", "What customers see and touch. Built with React, runs in the browser."),
        ("Kitchen → Backend", "The brain on the server. Built with PHP — receives requests, decides what to do."),
        ("Pantry → Database", "Stores the menu, every order and reservation. MySQL."),
        ("Waiter → Recommender", "A small helper that has 'seen' past orders and knows what pairs well."),
    ],
    index=3,
)

s_rows(
    "Chapter 3 · the big idea", "What is Software Construction?",
    "Construction is writing the actual code. SWEBOK lists what good construction looks "
    "like — five fundamentals. They are the spine of this deck.",
    [
        ("Minimizing Complexity", "keep it simple, so a human can understand it."),
        ("Anticipating Change", "build so tomorrow's change is easy, not painful."),
        ("Constructing for Verification", "build so it's easy to check that it works."),
        ("Reuse", "make a part once, use it everywhere."),
        ("Standards", "everyone follows the same rulebook."),
    ],
    index=4,
)

s_content(
    "Fundamental 01 / 05", "Minimizing Complexity",
    "The enemy is tangled code. The goal: make each piece small and obvious.",
    [
        ("The analogy", "One tidy serving hatch — not fifty secret doors into the kitchen."),
        ("In the code", "Every request to the server goes through ONE shared helper (api.ts)."),
        ("Why it pays off", "If the page-to-server link breaks, there is one place to look — not fifty."),
        ("Say this", "“Simple code is code a stranger can read at midnight and still understand.”"),
    ],
    index=5, wm="1",
)

s_content(
    "Fundamental 02 / 05", "Anticipating Change",
    "Software is never finished. Good construction makes change painless.",
    [
        ("Migrations", "Database changes apply (up) AND undo themselves (down) — like a save point."),
        ("Settings, not hard-coding", "The database address is read from the environment."),
        ("Move servers easily", "Change one setting, touch no code — a lamp with a plug, not wired into the wall."),
        ("Result", "Tomorrow's change is a small, safe, reversible step."),
    ],
    index=6, wm="2",
)

s_content(
    "Fundamental 03 / 05", "Constructing for Verification",
    "Verification means proving the software is correct, through testing — made easy on purpose.",
    [
        ("Swappable database", "Tests slot in a fake, throwaway database (setConnection) — fast and safe."),
        ("One command", "composer test runs every unit test; composer lint checks style."),
        ("Automated re-check", "Jenkins re-runs everything on each change, even calling /api/menu live."),
        ("Result", "Mistakes are caught before customers ever see them."),
    ],
    index=7, wm="3",
)

s_content(
    "Fundamental 04 / 05", "Reuse",
    "Build a part once, reuse it — less work, fewer bugs, consistent results.",
    [
        ("Design tokens", "Every colour and corner is defined once. Change the brand colour once — the whole site updates."),
        ("Shared API helper", "One 'api' object feeds the menu, basket, recommendations and reservations."),
        ("The analogy", "A house-blend sauce mixed once and used across many dishes."),
        ("Result", "Consistency for free, and half the code to maintain."),
    ],
    index=8, wm="4",
)

s_content(
    "Fundamental 05 / 05", "Standards in Construction",
    "Agreed rules, so the code reads like one careful author wrote it.",
    [
        ("PSR-12", "An automatic style checker keeps spacing and naming consistent in every file."),
        ("Strict types", "Every PHP file turns on strict mode — mistakes are caught early, not hidden."),
        ("One config", ".editorconfig gives every editor the same indentation and line endings."),
        ("Result", "A codebase that looks intentional, not accidental."),
    ],
    index=9, wm="5",
)

s_rows(
    "Part 2 · the toolchain", "The tools that ship it: DevOps",
    "Writing good code is cooking at home. DevOps turns it into a restaurant that serves "
    "the same dish, every day, to everyone.",
    [
        ("Git", "the time machine — records every change."),
        ("GitHub", "the shared home — the project online, one source of truth."),
        ("Docker", "the lunchbox — packs the app + all it needs to run anywhere."),
        ("Jenkins", "the robot inspector — checks, tests and builds on every change."),
    ],
    index=10,
)

s_content(
    "Tools 01 & 02 · version control + collaboration", "Git & GitHub",
    "Git remembers; GitHub shares.",
    [
        ("Git — commits", "A snapshot each time you finish work — like save points in a game. Nothing is ever lost."),
        ("Git — tidy history", "Commits are labelled (feat:, chore:) so the project's story reads clearly."),
        ("GitHub — the home", "Puts the history online: one source of truth for the team and the tools."),
        ("Why it matters", "Jenkins and the host watch GitHub and react when new code arrives."),
    ],
    index=11,
)

s_content(
    "Tool 03 · containerization", "Docker — the lunchbox",
    "Ends the “it works on my computer” problem.",
    [
        ("The container", "Packs the app AND everything it needs — exact PHP, extensions, settings — in one sealed box."),
        ("Same everywhere", "Open the box on any computer and it runs identically: laptop or cloud."),
        ("One command", "docker-compose starts the whole restaurant — PHP, MySQL, web server — together."),
        ("Your claim — correct", "“Docker containerizes the app so it moves local → cloud without breaking.”"),
    ],
    index=12,
)

s_content(
    "Tool 04 · automation (CI/CD)", "Jenkins — the robot inspector",
    "On every change, Jenkins runs a fixed checklist automatically. Fail a step — the code can't ship.",
    [
        ("The pipeline", "Checkout → install → lint & style → unit tests → build image → smoke test → deploy."),
        ("Testing & shipping", "That is Jenkins' real job — proving the code works, then releasing it."),
        ("On migrations", "It runs a migration ONLY to build a throwaway test database for the smoke test."),
        ("Who owns migrations", "Phinx does — Jenkins just presses the button during integration testing."),
    ],
    index=13,
)

closing()

# ── save ──────────────────────────────────────────────────────────────────────
out_dir = Path(__file__).resolve().parent / "assets"
out_dir.mkdir(parents=True, exist_ok=True)
out_path = out_dir / "Maison_Cafu_Presentation.pptx"
prs.save(str(out_path))
print(f"Saved {len(prs.slides)} slides -> {out_path}")
